#!/usr/bin/env python3
"""
XOYO Office Agent — Native PPTX/DOCX generation.
Port: 8056
"""
import os
import re
import subprocess
import time
import logging
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel, Field
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from docx import Document
    OFFICE_LIBS_AVAILABLE = True
except ImportError:
    OFFICE_LIBS_AVAILABLE = False
    # Mock classes for type hints
    class Presentation: pass
    class Document: pass
import uvicorn

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(title="XOYO Office Agent")

OUTPUT_DIR = os.path.expanduser("~/xoyo/output/documents/")
os.makedirs(OUTPUT_DIR, exist_ok=True)


def _sanitize_filename(name: str) -> str:
    """Remove dangerous characters from filenames to prevent path traversal."""
    # Strip directory separators and null bytes
    name = name.replace("/", "").replace("\\", "").replace("\0", "")
    # Remove other unsafe chars
    name = re.sub(r'[<>:"|?*]', '', name)
    # Replace spaces with underscores
    name = name.replace(" ", "_")
    # Fallback if empty
    if not name.strip("._"):
        name = f"xoyo_doc_{int(time.time())}"
    return name[:100]


class SlideModel(BaseModel):
    title: str = ""
    content: str = ""


class PptxRequest(BaseModel):
    title: str = Field(..., min_length=1)
    slides: list[SlideModel] = Field(..., min_length=1)
    open_file: bool = True


class DocxRequest(BaseModel):
    title: str = Field(..., min_length=1)
    content: str = ""
    open_file: bool = True

class InspectRequest(BaseModel):
    filepath: str

class EditRequest(BaseModel):
    filepath: str
    target_text: str
    replacement_text: str


@app.post("/create_pptx")
def create_pptx(request: PptxRequest):
    if not OFFICE_LIBS_AVAILABLE:
        raise HTTPException(status_code=501, detail="Office libraries not installed")
    try:
        prs = Presentation()

        # Add a title slide first
        title_slide_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_slide_layout)
        if title_slide.shapes.title:
            title_slide.shapes.title.text = request.title

        for slide_data in request.slides:
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)

            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = slide_data.title

            if len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.text = slide_data.content

        safe_name = _sanitize_filename(request.title)
        filepath = os.path.join(OUTPUT_DIR, f"{safe_name}.pptx")
        prs.save(filepath)
        logger.info(f"PPTX saved: {filepath}")

        if request.open_file:
            subprocess.Popen(
                ["xdg-open", filepath],
                stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
                start_new_session=True
            )

        return {"status": "ok", "filepath": filepath}
    except Exception as e:
        logger.error(f"PPTX creation failed: {e}")
        raise HTTPException(status_code=500, detail=f"PPTX creation failed: {type(e).__name__}")


@app.post("/create_docx")
def create_docx(request: DocxRequest):
    if not OFFICE_LIBS_AVAILABLE:
        raise HTTPException(status_code=501, detail="Office libraries not installed")
    try:
        document = Document()
        document.add_heading(request.title, 0)

        # Split content into paragraphs for better formatting
        paragraphs = request.content.split("\n") if request.content else [""]
        for para in paragraphs:
            if para.strip():
                document.add_paragraph(para.strip())

        safe_name = _sanitize_filename(request.title)
        filepath = os.path.join(OUTPUT_DIR, f"{safe_name}.docx")
        document.save(filepath)
        logger.info(f"DOCX saved: {filepath}")

        if request.open_file:
            subprocess.Popen(
                ["xdg-open", filepath],
                stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
                start_new_session=True
            )

        return {"status": "ok", "filepath": filepath}
    except Exception as e:
        logger.error(f"DOCX creation failed: {e}")
        raise HTTPException(status_code=500, detail=f"DOCX creation failed: {type(e).__name__}")

@app.post("/inspect_docx")
def inspect_docx(request: InspectRequest):
    if not OFFICE_LIBS_AVAILABLE:
        raise HTTPException(status_code=501, detail="Office libraries not installed")
    try:
        if not os.path.exists(request.filepath):
            raise HTTPException(status_code=404, detail="File not found")
        document = Document(request.filepath)
        paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
        return {"status": "ok", "paragraphs": paragraphs}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/edit_docx")
def edit_docx(request: EditRequest):
    if not OFFICE_LIBS_AVAILABLE:
        raise HTTPException(status_code=501, detail="Office libraries not installed")
    try:
        if not os.path.exists(request.filepath):
            raise HTTPException(status_code=404, detail="File not found")
        document = Document(request.filepath)
        modified = False
        for para in document.paragraphs:
            if request.target_text in para.text:
                para.text = para.text.replace(request.target_text, request.replacement_text)
                modified = True
        if not modified:
            return {"status": "unchanged", "message": "Target text not found"}
        document.save(request.filepath)
        return {"status": "ok", "message": "Document updated"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/inspect_pptx")
def inspect_pptx(request: InspectRequest):
    if not OFFICE_LIBS_AVAILABLE:
        raise HTTPException(status_code=501, detail="Office libraries not installed")
    try:
        if not os.path.exists(request.filepath):
            raise HTTPException(status_code=404, detail="File not found")
        prs = Presentation(request.filepath)
        slides_data = []
        for i, slide in enumerate(prs.slides):
            text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_runs.append(shape.text.strip())
            slides_data.append({"slide_index": i, "text": text_runs})
        return {"status": "ok", "slides": slides_data}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/edit_pptx")
def edit_pptx(request: EditRequest):
    if not OFFICE_LIBS_AVAILABLE:
        raise HTTPException(status_code=501, detail="Office libraries not installed")
    try:
        if not os.path.exists(request.filepath):
            raise HTTPException(status_code=404, detail="File not found")
        prs = Presentation(request.filepath)
        modified = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and request.target_text in shape.text:
                    shape.text = shape.text.replace(request.target_text, request.replacement_text)
                    modified = True
        if not modified:
            return {"status": "unchanged", "message": "Target text not found"}
        prs.save(request.filepath)
        return {"status": "ok", "message": "Presentation updated"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/health")
def health():
    return {"status": "ok", "service": "office_agent", "port": 8056}


if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8056)
